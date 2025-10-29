"""
Document Processor for Mental Health Reports
Extracts and analyzes text from various file formats including PDFs
"""

import os
import logging
from typing import Dict, Optional, Tuple
import tempfile

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process various document formats for mental health analysis"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._extract_pdf,
            '.txt': self._extract_txt,
            '.docx': self._extract_docx,
            '.doc': self._extract_docx,
            '.xlsx': self._extract_xlsx,
            '.xls': self._extract_xlsx,
            '.pptx': self._extract_pptx,
            '.html': self._extract_html,
            '.htm': self._extract_html,
            '.md': self._extract_txt,
        }
        
    def is_supported(self, filename: str) -> bool:
        """Check if file format is supported"""
        ext = os.path.splitext(filename.lower())[1]
        return ext in self.supported_formats
    
    def extract_text(self, file_path: str) -> Tuple[str, Dict]:
        """
        Extract text from document
        Returns: (extracted_text, metadata)
        """
        try:
            ext = os.path.splitext(file_path.lower())[1]
            
            if ext not in self.supported_formats:
                return "", {"error": f"Unsupported format: {ext}"}
            
            extractor = self.supported_formats[ext]
            text, metadata = extractor(file_path)
            
            logger.info(f"✅ Extracted {len(text)} characters from {ext} file")
            return text, metadata
            
        except Exception as e:
            logger.error(f"❌ Error extracting text: {e}")
            return "", {"error": str(e)}
    
    def _extract_pdf(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from PDF using multiple methods"""
        text = ""
        metadata = {"format": "pdf", "pages": 0}
        
        try:
            # Try pdfplumber first (better for complex PDFs)
            import pdfplumber
            
            with pdfplumber.open(file_path) as pdf:
                metadata["pages"] = len(pdf.pages)
                
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n\n--- Page {page_num} ---\n\n{page_text}"
                
                # Extract metadata
                if pdf.metadata:
                    metadata.update({
                        "title": pdf.metadata.get("Title", ""),
                        "author": pdf.metadata.get("Author", ""),
                        "subject": pdf.metadata.get("Subject", ""),
                    })
            
            logger.info(f"✅ Extracted PDF with pdfplumber: {metadata['pages']} pages")
            
        except Exception as e1:
            logger.warning(f"⚠️ pdfplumber failed: {e1}, trying PyPDF2...")
            
            try:
                # Fallback to PyPDF2
                import PyPDF2
                
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    metadata["pages"] = len(pdf_reader.pages)
                    
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        page_text = page.extract_text()
                        if page_text:
                            text += f"\n\n--- Page {page_num} ---\n\n{page_text}"
                    
                    # Extract metadata
                    if pdf_reader.metadata:
                        metadata.update({
                            "title": pdf_reader.metadata.get("/Title", ""),
                            "author": pdf_reader.metadata.get("/Author", ""),
                            "subject": pdf_reader.metadata.get("/Subject", ""),
                        })
                
                logger.info(f"✅ Extracted PDF with PyPDF2: {metadata['pages']} pages")
                
            except Exception as e2:
                logger.error(f"❌ Both PDF extractors failed: {e2}")
                raise Exception(f"Failed to extract PDF: {e2}")
        
        return text.strip(), metadata
    
    def _extract_txt(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            metadata = {
                "format": "text",
                "size": len(text),
                "lines": text.count('\n') + 1
            }
            
            return text, metadata
            
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    
                    metadata = {
                        "format": "text",
                        "encoding": encoding,
                        "size": len(text)
                    }
                    
                    return text, metadata
                except:
                    continue
            
            raise Exception("Failed to decode text file")
    
    def _extract_docx(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from Word document"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            
            # Extract paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            text = '\n\n'.join(paragraphs)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells)
                    text += f"\n{row_text}"
            
            metadata = {
                "format": "docx",
                "paragraphs": len(paragraphs),
                "tables": len(doc.tables)
            }
            
            # Extract core properties
            if doc.core_properties:
                metadata.update({
                    "title": doc.core_properties.title or "",
                    "author": doc.core_properties.author or "",
                    "subject": doc.core_properties.subject or "",
                })
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"❌ Error extracting DOCX: {e}")
            raise
    
    def _extract_xlsx(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from Excel spreadsheet"""
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, data_only=True)
            text = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n\n=== Sheet: {sheet_name} ===\n\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text += f"{row_text}\n"
            
            metadata = {
                "format": "xlsx",
                "sheets": len(wb.sheetnames),
                "sheet_names": wb.sheetnames
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"❌ Error extracting XLSX: {e}")
            raise
    
    def _extract_pptx(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from PowerPoint presentation"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n\n=== Slide {slide_num} ===\n\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += f"{shape.text}\n"
            
            metadata = {
                "format": "pptx",
                "slides": len(prs.slides)
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"❌ Error extracting PPTX: {e}")
            raise
    
    def _extract_html(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'lxml')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            metadata = {
                "format": "html",
                "title": soup.title.string if soup.title else ""
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"❌ Error extracting HTML: {e}")
            raise
    
    def analyze_mental_health_report(self, text: str, metadata: Dict) -> Dict:
        """
        Analyze extracted text for mental health report indicators
        Returns structured analysis
        """
        analysis = {
            "document_type": "unknown",
            "key_sections": [],
            "indicators": [],
            "summary": ""
        }
        
        text_lower = text.lower()
        
        # Detect document type
        if any(term in text_lower for term in ['diagnosis', 'patient', 'clinical', 'assessment', 'treatment']):
            analysis["document_type"] = "clinical_report"
        elif any(term in text_lower for term in ['therapy', 'session', 'counseling', 'progress']):
            analysis["document_type"] = "therapy_notes"
        elif any(term in text_lower for term in ['test', 'score', 'evaluation', 'questionnaire']):
            analysis["document_type"] = "assessment"
        elif any(term in text_lower for term in ['prescription', 'medication', 'dosage']):
            analysis["document_type"] = "prescription"
        
        # Identify key sections
        section_keywords = {
            "diagnosis": ["diagnosis", "diagnosed with", "condition"],
            "symptoms": ["symptoms", "experiencing", "reports"],
            "treatment": ["treatment", "therapy", "medication", "intervention"],
            "progress": ["progress", "improvement", "response to treatment"],
            "recommendations": ["recommend", "suggestion", "next steps"]
        }
        
        for section, keywords in section_keywords.items():
            if any(keyword in text_lower for keyword in keywords):
                analysis["key_sections"].append(section)
        
        # Detect mental health indicators
        indicators = {
            "depression": ["depression", "depressed", "low mood", "anhedonia"],
            "anxiety": ["anxiety", "anxious", "panic", "worry", "fear"],
            "trauma": ["trauma", "ptsd", "flashback", "traumatic"],
            "substance": ["substance", "alcohol", "drug", "addiction"],
            "suicidal": ["suicidal", "self-harm", "suicide", "ideation"],
            "bipolar": ["bipolar", "manic", "mania", "mood swings"],
            "psychosis": ["psychosis", "hallucination", "delusion"],
            "eating": ["eating disorder", "anorexia", "bulimia", "binge"]
        }
        
        for indicator, keywords in indicators.items():
            if any(keyword in text_lower for keyword in keywords):
                analysis["indicators"].append(indicator)
        
        # Generate summary
        analysis["summary"] = f"Document appears to be a {analysis['document_type']} "
        if analysis["indicators"]:
            analysis["summary"] += f"with indicators for: {', '.join(analysis['indicators'])}. "
        if analysis["key_sections"]:
            analysis["summary"] += f"Contains sections on: {', '.join(analysis['key_sections'])}."
        
        return analysis


# Global instance
document_processor = DocumentProcessor()
